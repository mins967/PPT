"""
PARIS Cultural Presentation — PPT Generator
Generates a .pptx file matching the HTML framework layout.
Images are downloaded from Unsplash and embedded.
Animations should be added manually in PowerPoint afterward.
"""

import os
import urllib.request
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Constants ──
SLIDE_W = Emu(12192000)  # 13.333"
SLIDE_H = Emu(6858000)   # 7.5"

BLUE   = RGBColor(0x00, 0x26, 0x54)
RED    = RGBColor(0xCE, 0x11, 0x26)
BG     = RGBColor(0xF4, 0xF3, 0xF1)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLUE_A = RGBColor(0x00, 0x1A, 0x3A)  # darker blue for ghost text

FONT_TITLE = 'Bebas Neue'
FONT_BODY  = 'DM Sans'
FONT_SERIF = 'Playfair Display'

IMAGES = {
    'eiffel':     'https://images.unsplash.com/photo-1502602898657-3e91760cbb34?w=1200&q=85',
    'seine':      'https://images.unsplash.com/photo-1499856871958-5b9627545d1a?w=1200&q=85',
    'art':        'https://images.unsplash.com/photo-1509439581779-6298f75bf6e5?w=1200&q=85',
    'eiffel_old': 'https://images.unsplash.com/photo-1431274172761-fca41d930114?w=1200&q=85',
    'modern':     'https://images.unsplash.com/photo-1563804447971-6e113ab80713?w=1400&q=85',
    'bastille':   'https://images.unsplash.com/photo-1514190051997-0f6f39ca5cde?w=1200&q=85',
    'crepe':      'https://images.unsplash.com/photo-1741202364284-d7f0da126f63?w=1200&q=85',
    'cuisine':    'https://images.unsplash.com/photo-1555396273-367ea4eb4db5?w=900&q=85',
}


def pct_w(p):
    return int(SLIDE_W.emu * p / 100)

def pct_h(p):
    return int(SLIDE_H.emu * p / 100)


def download_image(url):
    try:
        req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
        with urllib.request.urlopen(req, timeout=15) as resp:
            return BytesIO(resp.read())
    except Exception as e:
        print(f"  [!] Image download failed: {e}")
        return None


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width or Pt(1)
    return shape


def add_text_box(slide, left, top, width, height, text, font_name=FONT_BODY,
                 font_size=Pt(12), font_color=WHITE, bold=False, italic=False,
                 alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = font_size
        run.font.color.rgb = font_color
        run.font.bold = bold
        run.font.italic = italic

        p.alignment = alignment
        if spacing is not None:
            p.space_before = spacing
            p.space_after = spacing
    return txBox


def add_image_safe(slide, img_stream, left, top, width, height):
    if img_stream:
        img_stream.seek(0)
        return slide.shapes.add_picture(img_stream, left, top, width, height)
    else:
        return add_rect(slide, left, top, width, height, fill_color=BLUE_A)


# ── Image Cache ──
print("Downloading images...")
img_cache = {}
for key, url in IMAGES.items():
    print(f"  [{key}] ", end="")
    img_cache[key] = download_image(url)
    print("OK" if img_cache[key] else "FAILED (placeholder will be used)")

# ── Create Presentation ──
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]

W = SLIDE_W.emu
H = SLIDE_H.emu


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 1: HERO
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
print("Building Slide 1: HERO...")
s1 = prs.slides.add_slide(blank_layout)
set_slide_bg(s1, BLUE)

left_w = pct_w(52)
right_w = W - left_w

# Red vertical stripe (left edge)
add_rect(s1, 0, 0, Emu(36000), H, fill_color=RED)

# Top bar
add_text_box(s1, Emu(330000), Emu(200000), Emu(2500000), Emu(200000),
             "CULTURAL PRESENTATION · 2025",
             font_size=Pt(7), font_color=RGBColor(0x66, 0x77, 0x88))

add_rect(s1, Emu(left_w - 1300000), Emu(180000), Emu(1200000), Emu(230000), fill_color=RED)
add_text_box(s1, Emu(left_w - 1280000), Emu(195000), Emu(1160000), Emu(200000),
             "FRANCE · EUROPE",
             font_size=Pt(7), font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# PARIS title
add_text_box(s1, Emu(330000), Emu(1800000), Emu(left_w - 500000), Emu(2800000),
             "PARIS",
             font_name=FONT_TITLE, font_size=Pt(140), font_color=WHITE,
             bold=True, anchor=MSO_ANCHOR.MIDDLE)

# Subtitle
add_text_box(s1, Emu(330000), Emu(4300000), Emu(left_w - 500000), Emu(300000),
             "HERITAGE · ART · GASTRONOMY · FASHION · HISTORY",
             font_size=Pt(8), font_color=RGBColor(0x66, 0x77, 0x88))

# Stats row
stat_w = left_w // 3
stat_y = Emu(5400000)
stat_data = [("2,000+", "YEARS OF HISTORY"), ("2.1M", "CITY POPULATION"), ("#1", "MOST VISITED CITY")]
add_rect(s1, 0, stat_y - Emu(30000), left_w, Emu(1), fill_color=RGBColor(0x11, 0x33, 0x55))

for i, (val, lbl) in enumerate(stat_data):
    x = Emu(stat_w * i + 250000)
    add_text_box(s1, x, stat_y, Emu(stat_w - 200000), Emu(320000),
                 val, font_name=FONT_TITLE, font_size=Pt(26), font_color=WHITE)
    add_text_box(s1, x, stat_y + Emu(350000), Emu(stat_w - 200000), Emu(200000),
                 lbl, font_size=Pt(7), font_color=RGBColor(0x55, 0x66, 0x77))

# Flag strip
flag_y = H - Emu(46000)
flag_h = Emu(46000)
flag_w = left_w // 3
add_rect(s1, 0, flag_y, flag_w, flag_h, fill_color=BLUE)
add_rect(s1, flag_w, flag_y, flag_w, flag_h, fill_color=WHITE)
add_rect(s1, flag_w * 2, flag_y, flag_w, flag_h, fill_color=RED)

# Right: main image (top 65%)
img_top_h = int(H * 0.65)
add_image_safe(s1, img_cache['eiffel'], left_w, 0, right_w, img_top_h)

# Photo label
add_text_box(s1, left_w + Emu(200000), img_top_h - Emu(350000), Emu(2000000), Emu(250000),
             "TOUR EIFFEL", font_name=FONT_TITLE, font_size=Pt(16), font_color=WHITE)
add_text_box(s1, left_w + Emu(200000) + Emu(2100000), img_top_h - Emu(310000), Emu(2500000), Emu(200000),
             "CHAMP DE MARS · PARIS", font_size=Pt(7), font_color=RGBColor(0x99, 0xAA, 0xBB),
             alignment=PP_ALIGN.RIGHT)

# Right: bottom images (35%)
img_bot_h = H - img_top_h
sub_w = right_w // 2
add_image_safe(s1, img_cache['seine'], left_w, img_top_h, sub_w, img_bot_h)
add_image_safe(s1, img_cache['cuisine'], left_w + sub_w, img_top_h, sub_w, img_bot_h)

add_text_box(s1, left_w + Emu(140000), img_top_h + img_bot_h - Emu(280000),
             Emu(2000000), Emu(200000),
             "SEINE · ÎLE DE LA CITÉ", font_size=Pt(7), font_color=RGBColor(0xAA, 0xAA, 0xAA))
add_text_box(s1, left_w + sub_w + Emu(140000), img_top_h + img_bot_h - Emu(280000),
             Emu(2000000), Emu(200000),
             "PARISIAN CUISINE", font_size=Pt(7), font_color=RGBColor(0xAA, 0xAA, 0xAA))

# Red bottom line
add_rect(s1, 0, H - Emu(36000), W, Emu(36000), fill_color=RED)

# Keyword tags on right panel
kw_x = left_w + Emu(180000)
for i, kw in enumerate(["CITY OF LIGHT", "CAPITAL OF FRANCE", "48°51′N · 2°21′E"]):
    tag = add_rect(s1, kw_x, Emu(200000 + i * 250000), Emu(1600000), Emu(200000),
                   fill_color=RGBColor(0x00, 0x20, 0x44))
    add_rect(s1, kw_x, Emu(200000 + i * 250000), Emu(20000), Emu(200000), fill_color=RED)
    add_text_box(s1, kw_x + Emu(50000), Emu(210000 + i * 250000),
                 Emu(1500000), Emu(180000),
                 kw, font_size=Pt(7), font_color=WHITE)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 2: LOCATION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
print("Building Slide 2: LOCATION...")
s2 = prs.slides.add_slide(blank_layout)
set_slide_bg(s2, BG)

left_w = pct_w(55)
right_w = W - left_w

add_image_safe(s2, img_cache['seine'], 0, 0, left_w, H)

label_bg = add_rect(s2, Emu(290000), H - Emu(500000), Emu(1400000), Emu(260000), fill_color=BLUE)
add_text_box(s2, Emu(340000), H - Emu(475000), Emu(1300000), Emu(210000),
             "ÎLE-DE-FRANCE", font_size=Pt(8), font_color=WHITE, bold=True)

add_rect(s2, left_w, 0, right_w, H, fill_color=BLUE)

# Ghost number
add_text_box(s2, Emu(left_w + right_w - 1800000), Emu(-100000), Emu(2000000), Emu(1400000),
             "02", font_name=FONT_TITLE, font_size=Pt(120), font_color=RGBColor(0x00, 0x2E, 0x5C),
             alignment=PP_ALIGN.RIGHT)

tx = left_w + Emu(470000)

add_text_box(s2, tx, Emu(2200000), Emu(right_w - 900000), Emu(200000),
             "LOCATION", font_size=Pt(8), font_color=RED, bold=True)

add_text_box(s2, tx, Emu(2500000), Emu(right_w - 900000), Emu(1000000),
             "WHERE\nIS IT?", font_name=FONT_TITLE, font_size=Pt(54), font_color=WHITE)

add_rect(s2, tx, Emu(3600000), Emu(440000), Emu(18000), fill_color=RED)

kw_y = Emu(3800000)
keywords = ["NORTH-CENTRAL FRANCE", "SEINE RIVER", "48°N · 2°E", "CAPITAL CITY", "WESTERN EUROPE"]
for i, kw in enumerate(keywords):
    dot = add_rect(s2, tx, kw_y + Emu(i * 230000) + Emu(50000),
                   Emu(36000), Emu(36000), fill_color=RED)
    add_text_box(s2, tx + Emu(100000), kw_y + Emu(i * 230000),
                 Emu(right_w - 1100000), Emu(200000),
                 kw, font_size=Pt(10), font_color=RGBColor(0x88, 0x99, 0xAA))


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 3: WHY PARIS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
print("Building Slide 3: WHY PARIS...")
s3 = prs.slides.add_slide(blank_layout)
set_slide_bg(s3, BG)

left_w = pct_w(42)
right_w = W - left_w

add_image_safe(s3, img_cache['art'], 0, 0, left_w, H)

badge = add_rect(s3, Emu(330000), Emu(330000), Emu(1100000), Emu(240000), fill_color=RED)
add_text_box(s3, Emu(370000), Emu(345000), Emu(1020000), Emu(210000),
             "WHY PARIS", font_size=Pt(7), font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(s3, left_w, 0, right_w, H, fill_color=BLUE)

tx = left_w + Emu(470000)

add_text_box(s3, tx, Emu(430000), Emu(right_w - 900000), Emu(200000),
             "WHY I LOVE IT", font_size=Pt(8), font_color=RED, bold=True)

add_text_box(s3, tx, Emu(700000), Emu(right_w - 700000), Emu(1600000),
             "REASONS\nTO LOVE\nPARIS", font_name=FONT_TITLE, font_size=Pt(52), font_color=WHITE)

add_rect(s3, tx, Emu(2450000), Emu(440000), Emu(18000), fill_color=RED)

reasons = [
    ("01", "ART & ARCHITECTURE", "Louvre · Notre-Dame · Haussmann"),
    ("02", "GASTRONOMY & CAFÉ", "Bistros · Patisseries · Michelin"),
    ("03", "FASHION & STYLE", "Haute Couture · Chanel · Savoir-Faire"),
]

ry = Emu(2700000)
for i, (num, kw, sub) in enumerate(reasons):
    y_off = ry + Emu(i * 1100000)
    add_text_box(s3, tx, y_off, Emu(400000), Emu(350000),
                 num, font_name=FONT_TITLE, font_size=Pt(28), font_color=RED)
    add_text_box(s3, tx + Emu(450000), y_off, Emu(right_w - 1400000), Emu(300000),
                 kw, font_name=FONT_TITLE, font_size=Pt(22), font_color=WHITE)
    add_text_box(s3, tx + Emu(450000), y_off + Emu(320000), Emu(right_w - 1400000), Emu(200000),
                 sub, font_size=Pt(8), font_color=RGBColor(0x66, 0x77, 0x88))
    if i < 2:
        add_rect(s3, tx, y_off + Emu(750000), Emu(right_w - 950000), Emu(8000),
                 fill_color=RGBColor(0x00, 0x30, 0x60))

# Ghost letter
add_text_box(s3, Emu(W - 2500000), H - Emu(2200000), Emu(2500000), Emu(2200000),
             "P", font_name=FONT_TITLE, font_size=Pt(200), font_color=RGBColor(0x00, 0x2A, 0x58),
             alignment=PP_ALIGN.RIGHT)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 4: HISTORICAL FUN FACT
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
print("Building Slide 4: HISTORICAL FUN FACT...")
s4 = prs.slides.add_slide(blank_layout)
set_slide_bg(s4, BLUE)

half_w = pct_w(50)

# Ghost year
add_text_box(s4, Emu(-150000), H - Emu(1700000), Emu(5000000), Emu(1700000),
             "1889", font_name=FONT_TITLE, font_size=Pt(150), font_color=RGBColor(0x00, 0x2E, 0x5C))

tx = Emu(510000)
add_text_box(s4, tx, Emu(1800000), Emu(half_w - 800000), Emu(200000),
             "HISTORICAL FUN FACT", font_size=Pt(8), font_color=RED, bold=True)

add_text_box(s4, tx, Emu(2150000), Emu(half_w - 800000), Emu(1800000),
             "Built to\nbe torn\ndown",
             font_name=FONT_SERIF, font_size=Pt(44), font_color=WHITE, italic=True)

# Pills
pill_y = Emu(4200000)
pills = [("EIFFEL TOWER", True), ("1889", False), ("TEMPORARY", False), ("IRON GIANT", False), ("300M TALL", False)]
px = tx
for label_text, accent in pills:
    pw = Emu(len(label_text) * 72000 + 220000)
    fill = RED if accent else None
    border = None if accent else RGBColor(0x44, 0x55, 0x66)
    pill_shape = add_rect(s4, px, pill_y, pw, Emu(240000),
                          fill_color=fill, border_color=border, border_width=Pt(0.75))
    fc = WHITE if accent else RGBColor(0xAA, 0xBB, 0xCC)
    add_text_box(s4, px + Emu(20000), pill_y + Emu(20000), pw - Emu(40000), Emu(200000),
                 label_text, font_size=Pt(8), font_color=fc, alignment=PP_ALIGN.CENTER)
    px += pw + Emu(70000)

# Right: image
add_image_safe(s4, img_cache['eiffel_old'], half_w, 0, half_w, H)

# Gradient overlay (semi-transparent rect)
add_rect(s4, half_w, 0, int(half_w * 0.5), H, fill_color=BLUE)

# Stamp
stamp_x = W - Emu(1700000)
stamp_y = H - Emu(1300000)
stamp = add_rect(s4, stamp_x, stamp_y, Emu(1300000), Emu(900000),
                 border_color=RGBColor(0x66, 0x77, 0x88), border_width=Pt(1.5))
add_text_box(s4, stamp_x + Emu(100000), stamp_y + Emu(100000), Emu(1100000), Emu(500000),
             "1889", font_name=FONT_TITLE, font_size=Pt(40), font_color=WHITE,
             alignment=PP_ALIGN.CENTER)
add_text_box(s4, stamp_x + Emu(100000), stamp_y + Emu(580000), Emu(1100000), Emu(200000),
             "WORLD'S FAIR · PARIS", font_size=Pt(7), font_color=RGBColor(0x88, 0x99, 0xAA),
             alignment=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 5: MODERN FUN FACT
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
print("Building Slide 5: MODERN FUN FACT...")
s5 = prs.slides.add_slide(blank_layout)
set_slide_bg(s5, BG)

top_h = pct_h(60)
bot_h = H - top_h

add_image_safe(s5, img_cache['modern'], 0, 0, W, top_h)

badge5 = add_rect(s5, Emu(440000), Emu(330000), Emu(1600000), Emu(260000), fill_color=RED)
add_text_box(s5, Emu(480000), Emu(345000), Emu(1520000), Emu(230000),
             "MODERN FUN FACT", font_size=Pt(8), font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(s5, 0, top_h, W, Emu(8000), fill_color=RGBColor(0xDD, 0xDC, 0xD9))

items = [
    ("#1", "TOURISM", "Most Visited\nCity · World"),
    ("100M+", "VISITORS / YEAR", "Pre-Pandemic\nRecord"),
    ("4,500", "RESTAURANTS", "Michelin Star\nCapital"),
    ("132", "MUSEUMS", "City of\nLight · Art"),
]

item_w = W // 4
for i, (num, lbl, kw) in enumerate(items):
    ix = item_w * i
    iy = top_h + Emu(290000)

    add_text_box(s5, ix + Emu(330000), iy, Emu(item_w - 400000), Emu(400000),
                 num, font_name=FONT_TITLE, font_size=Pt(30), font_color=BLUE)
    add_text_box(s5, ix + Emu(330000), iy + Emu(380000), Emu(item_w - 400000), Emu(200000),
                 lbl, font_size=Pt(8), font_color=RED)
    add_text_box(s5, ix + Emu(330000), iy + Emu(600000), Emu(item_w - 400000), Emu(400000),
                 kw, font_size=Pt(8), font_color=RGBColor(0x77, 0x88, 0x99))

    if i < 3:
        add_rect(s5, ix + item_w - Emu(4000), top_h, Emu(8000), bot_h,
                 fill_color=RGBColor(0xE8, 0xE7, 0xE4))


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 6: CULTURAL TRADITION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
print("Building Slide 6: CULTURAL TRADITION...")
s6 = prs.slides.add_slide(blank_layout)
set_slide_bg(s6, WHITE)

left_w = pct_w(45)
right_w = W - left_w

add_rect(s6, 0, 0, left_w, H, fill_color=BLUE)

# Ghost letter
add_text_box(s6, Emu(left_w - 2500000), H - Emu(2500000), Emu(3000000), Emu(2500000),
             "F", font_name=FONT_TITLE, font_size=Pt(240), font_color=RGBColor(0x00, 0x2A, 0x58),
             alignment=PP_ALIGN.RIGHT)

tx = Emu(470000)
add_text_box(s6, tx, Emu(1800000), Emu(left_w - 800000), Emu(200000),
             "CULTURAL TRADITION", font_size=Pt(8), font_color=RED, bold=True)

add_text_box(s6, tx, Emu(2100000), Emu(left_w - 700000), Emu(2200000),
             "LA\nFÊTE\nNATIO-\nNALE",
             font_name=FONT_TITLE, font_size=Pt(60), font_color=WHITE)

tags = ["Bastille Day · July 14", "National Unity", "1789 Revolution", "Fireworks · Parade"]
ty = Emu(4800000)
for i, tag in enumerate(tags):
    add_rect(s6, tx, ty + Emu(i * 230000) + Emu(50000), Emu(220000), Emu(8000), fill_color=RED)
    add_text_box(s6, tx + Emu(290000), ty + Emu(i * 230000),
                 Emu(left_w - 1000000), Emu(200000),
                 tag.upper(), font_size=Pt(8), font_color=RGBColor(0x66, 0x77, 0x88))

# Right: top image
img_h = H // 2
add_image_safe(s6, img_cache['bastille'], left_w, 0, right_w, img_h)

# Right: bottom stats
stat_w = right_w // 2
stats = [("1789", "FOUNDED", "French Revolution"), ("233", "YEARS OLD", "Living Tradition")]
for i, (val, lbl, kw) in enumerate(stats):
    sx = left_w + stat_w * i
    sy = img_h
    add_rect(s6, sx, sy, Emu(8000) if i == 0 else 0, H - img_h, fill_color=RGBColor(0xEE, 0xED, 0xEB))

    add_text_box(s6, sx + Emu(290000), sy + Emu(250000), Emu(stat_w - 500000), Emu(500000),
                 val, font_name=FONT_TITLE, font_size=Pt(40), font_color=BLUE)
    add_text_box(s6, sx + Emu(290000), sy + Emu(750000), Emu(stat_w - 500000), Emu(200000),
                 lbl, font_size=Pt(8), font_color=RED)
    add_text_box(s6, sx + Emu(290000), sy + Emu(1000000), Emu(stat_w - 500000), Emu(200000),
                 kw.upper(), font_size=Pt(8), font_color=RGBColor(0x99, 0xAA, 0xBB))


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 7: FOOD — CRÊPE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
print("Building Slide 7: FOOD...")
s7 = prs.slides.add_slide(blank_layout)
set_slide_bg(s7, BLUE)

half_w = pct_w(50)

add_image_safe(s7, img_cache['crepe'], 0, 0, half_w, H)

badge7 = add_rect(s7, Emu(330000), Emu(330000), Emu(1300000), Emu(240000), fill_color=RED)
add_text_box(s7, Emu(370000), Emu(345000), Emu(1220000), Emu(210000),
             "CUISINE · ORIGINS", font_size=Pt(7), font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(s7, Emu(250000), H - Emu(350000), Emu(2500000), Emu(200000),
             "ICONIC FRENCH FOOD", font_size=Pt(7), font_color=RGBColor(0x66, 0x77, 0x88))

# Ghost letter
add_text_box(s7, Emu(W - 2200000), H - Emu(2000000), Emu(2200000), Emu(2000000),
             "C", font_name=FONT_TITLE, font_size=Pt(160), font_color=RGBColor(0x00, 0x2A, 0x58),
             alignment=PP_ALIGN.RIGHT)

tx = half_w + Emu(470000)

add_text_box(s7, tx, Emu(1200000), Emu(half_w - 800000), Emu(200000),
             "FOOD & CULTURE", font_size=Pt(8), font_color=RED, bold=True)

add_text_box(s7, tx, Emu(1500000), Emu(half_w - 700000), Emu(1200000),
             "CRÊPE", font_name=FONT_TITLE, font_size=Pt(72), font_color=WHITE)

add_rect(s7, tx, Emu(2900000), Emu(440000), Emu(24000), fill_color=RED)

facts = [
    ("ORIGIN", "Brittany, France · 13th Century"),
    ("BASE", "Buckwheat (Galette) · Wheat Flour"),
    ("TYPES", "Sweet (Dessert) · Savory (Galette)"),
    ("SYMBOL", "La Chandeleur · Feb 2nd"),
]

fy = Emu(3150000)
for i, (lbl, val) in enumerate(facts):
    y = fy + Emu(i * 280000)
    add_text_box(s7, tx, y, Emu(720000), Emu(200000),
                 lbl, font_size=Pt(7), font_color=RED, bold=True)
    add_text_box(s7, tx + Emu(780000), y, Emu(half_w - 1700000), Emu(200000),
                 val.upper(), font_size=Pt(10), font_color=RGBColor(0xCC, 0xCC, 0xCC))

# Stats
add_rect(s7, tx, Emu(4500000), Emu(half_w - 950000), Emu(8000), fill_color=RGBColor(0x11, 0x33, 0x55))

st_data = [("1B+", "EATEN PER YEAR IN FRANCE"), ("800+", "YEARS OF HISTORY")]
st_w = (half_w - 950000) // 2
for i, (sv, sl) in enumerate(st_data):
    sx = tx + st_w * i + (Emu(200000) if i == 1 else 0)
    add_text_box(s7, sx, Emu(4700000), Emu(st_w - 100000), Emu(400000),
                 sv, font_name=FONT_TITLE, font_size=Pt(30), font_color=WHITE)
    add_text_box(s7, sx, Emu(5100000), Emu(st_w - 100000), Emu(200000),
                 sl, font_size=Pt(7), font_color=RGBColor(0x55, 0x66, 0x77))


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 8: THANK YOU
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
print("Building Slide 8: THANK YOU...")
s8 = prs.slides.add_slide(blank_layout)
set_slide_bg(s8, BLUE)

# Red side lines
add_rect(s8, 0, 0, Emu(36000), H, fill_color=RED)
add_rect(s8, W - Emu(36000), 0, Emu(36000), H, fill_color=RED)

# Corner markers
corner_sz = Emu(260000)
corner_clr = RGBColor(0x22, 0x44, 0x66)
for (cx, cy) in [(Emu(330000), Emu(330000)),
                 (W - Emu(590000), Emu(330000)),
                 (Emu(330000), H - Emu(590000)),
                 (W - Emu(590000), H - Emu(590000))]:
    add_rect(s8, cx, cy, corner_sz, corner_sz,
             border_color=corner_clr, border_width=Pt(1))

# Ghost text
add_text_box(s8, Emu(1500000), Emu(1500000), Emu(W - 3000000), Emu(4000000),
             "MERCI", font_name=FONT_TITLE, font_size=Pt(280),
             font_color=RGBColor(0x00, 0x2A, 0x58), alignment=PP_ALIGN.CENTER)

# Center content
add_text_box(s8, Emu(2000000), Emu(2200000), Emu(W - 4000000), Emu(300000),
             "CULTURAL PRESENTATION · PARIS",
             font_size=Pt(8), font_color=RED, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(s8, Emu(2000000), Emu(2600000), Emu(W - 4000000), Emu(2000000),
             "THANK\nYOU", font_name=FONT_TITLE, font_size=Pt(120),
             font_color=WHITE, alignment=PP_ALIGN.CENTER)

add_rect(s8, (W - Emu(580000)) // 2, Emu(4700000), Emu(580000), Emu(24000), fill_color=RED)

add_text_box(s8, Emu(2000000), Emu(4900000), Emu(W - 4000000), Emu(300000),
             "HERITAGE · ART · GASTRONOMY · FASHION · HISTORY",
             font_size=Pt(8), font_color=RGBColor(0x55, 0x66, 0x77), alignment=PP_ALIGN.CENTER)

# Flag strip
flag_y8 = H - Emu(46000)
flag_w8 = W // 3
add_rect(s8, 0, flag_y8, flag_w8, Emu(46000), fill_color=BLUE)
add_rect(s8, flag_w8, flag_y8, flag_w8, Emu(46000), fill_color=WHITE)
add_rect(s8, flag_w8 * 2, flag_y8, flag_w8 + Emu(10000), Emu(46000), fill_color=RED)


# ── Save ──
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "PARIS_Presentation_v2.pptx")
prs.save(output_path)
print(f"\n✓ Saved: {output_path}")
print("  → Open in PowerPoint and add Morph transitions + entrance animations.")
